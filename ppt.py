import streamlit as st
import plotly.express as px
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import numpy as np
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches

# Load the data
data = pd.read_csv("C:/Users/pc click/immo_eliza_analysis/immoweb_data_cleaned.csv")

# Ensure data columns are correctly typed
st.sidebar.title("Data Overview and Filters")

# Display dataset overview
if st.sidebar.checkbox("Show Raw Data", value=False):
    st.subheader("Dataset Overview")
    st.write(data)

# Drop rows with missing values in key columns for analysis
filtered_data = data.dropna(subset=["Living_Area", "Price", "Number_of_Rooms"])

# Sidebar Filters
st.sidebar.subheader("Filter Options")
property_type = st.sidebar.multiselect(
    "Select Property Type:",
    options=data["Type_of_Property"].unique(),
    default=data["Type_of_Property"].unique()
)

state_of_building = st.sidebar.multiselect(
    "Select State of the Building:",
    options=data["State_of_the_Building"].unique(),
    default=data["State_of_the_Building"].unique()
)

# Apply filters
filtered_data = filtered_data[
    (filtered_data["Type_of_Property"].isin(property_type)) &
    (filtered_data["State_of_the_Building"].isin(state_of_building))
]

# Dashboard Title
st.title("Real Estate Interactive Dashboard")

# Display filtered data
st.header("Filtered Data")
st.write(f"Number of records after applying filters: {len(filtered_data)}")
st.dataframe(filtered_data)

# Key statistics (mean, median, etc.) for numerical columns
st.subheader("Key Statistics")
st.write(filtered_data.describe())

# Price vs Living Area Scatter Plot
st.header("Price vs Living Area")
scatter_fig = px.scatter(
    filtered_data,
    x="Living_Area",
    y="Price",
    color="Type_of_Property",
    size="Number_of_Rooms",
    hover_data=["Subtype_of_Property"],
    template="plotly_dark"
)

# Save the figure as an image
scatter_fig.write_image("price_vs_living_area.png")
st.plotly_chart(scatter_fig)

# Distribution of Prices
st.header("Price Distribution")
price_dist_fig = px.histogram(
    filtered_data,
    x="Price",
    nbins=30,
    color="Type_of_Property",
    title="Price Distribution by Property Type",
    template="plotly_dark"
)
price_dist_fig.write_image("price_distribution.png")
st.plotly_chart(price_dist_fig)

# Living Area Distribution
st.header("Living Area Distribution")
living_area_fig = px.histogram(
    filtered_data,
    x="Living_Area",
    nbins=30,
    color="Type_of_Property",
    title="Living Area Distribution by Property Type",
    template="plotly_dark"
)
living_area_fig.write_image("living_area_distribution.png")
st.plotly_chart(living_area_fig)

# Average Price per Locality
st.header("Average Price per Locality")
avg_price_locality = data.groupby("Municipality")["Price"].mean().reset_index().sort_values(by="Price", ascending=False)
locality_price_fig = px.bar(
    avg_price_locality.head(20),
    x="Municipality",
    y="Price",
    title="Top 20 Localities by Average Price",
    template="plotly_dark"
)
locality_price_fig.write_image("avg_price_locality.png")
st.plotly_chart(locality_price_fig)

# Boxplot of Price by Property Type
st.header("Price Distribution by Property Type")
boxplot_fig = px.box(
    filtered_data,
    x="Type_of_Property",
    y="Price",
    color="Type_of_Property",
    title="Boxplot of Price by Property Type",
    template="plotly_dark"
)
boxplot_fig.write_image("price_by_type.png")
st.plotly_chart(boxplot_fig)

# Pairwise Correlations
st.header("Pairwise Relationships")
pairwise_fig = px.scatter_matrix(
    filtered_data,
    dimensions=["Living_Area", "Price", "Number_of_Rooms"],
    color="Type_of_Property",
    title="Pairwise Correlation of Key Variables",
    template="plotly_dark"
)
pairwise_fig.write_image("pairwise_correlation.png")
st.plotly_chart(pairwise_fig)

# Factor Impact on Price (e.g., State of the Building, Number of Rooms, Fully Equipped Kitchen, etc.)
st.subheader("Impact of Features on Price")

# Filtered data based on selected features' impact
state_of_building_data = filtered_data[['State_of_the_Building', 'Price']].groupby('State_of_the_Building').mean().reset_index()
number_of_rooms_data = filtered_data[['Number_of_Rooms', 'Price']].groupby('Number_of_Rooms').mean().reset_index()
fully_equipped_kitchen_data = filtered_data[['Fully_Equipped_Kitchen', 'Price']].groupby('Fully_Equipped_Kitchen').mean().reset_index()
terrace_data = filtered_data[['Terrace', 'Price']].groupby('Terrace').mean().reset_index()

# State of the Building Impact
st.write("### Impact of State of the Building on Price")
st.write("How does the condition of the building affect its price?")
st.write(state_of_building_data)

# Number of Rooms Impact
st.write("### Impact of Number of Rooms on Price")
st.write("How does the number of rooms affect the price?")
st.write(number_of_rooms_data)

# Fully Equipped Kitchen Impact
st.write("### Impact of Fully Equipped Kitchen on Price")
st.write("How does having a fully equipped kitchen influence the price?")
st.write(fully_equipped_kitchen_data)

# Terrace Impact
st.write("### Impact of Terrace on Price")
st.write("Does having a terrace affect the price?")
st.write(terrace_data)

# Save transformed data if needed
filtered_data.to_csv("C:/Users/pc click/immo_eliza_analysis/immoweb_data_filtered.csv", index=False)
st.write("Filtered data saved to 'immoweb_data_filtered.csv'.")

# Creating PowerPoint from Saved Images
prs = Presentation()

# Add first slide with Price vs Living Area graph
slide_layout = prs.slide_layouts[5]  # Blank slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Price vs Living Area"
img_path = "price_vs_living_area.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5))

# Add second slide with Price Distribution graph
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Price Distribution"
img_path = "price_distribution.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5))

# Add third slide with Living Area Distribution graph
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Living Area Distribution"
img_path = "living_area_distribution.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5))

# Add fourth slide with Average Price per Locality graph
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Average Price per Locality"
img_path = "avg_price_locality.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5))

# Add fifth slide with Boxplot of Price by Property Type
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Boxplot of Price by Property Type"
img_path = "price_by_type.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5))

# Save the PowerPoint presentation
prs.save("Real_Estate_Analysis.pptx")
st.write("PowerPoint saved as 'Real_Estate_Analysis.pptx'.")

